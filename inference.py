import cv2
import os
import re
import numpy as np
from ultralytics import YOLO
import easyocr
import spacy
import pandas as pd
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import torch
import moondream as md

# --- INITIALIZATION (Load models only once) ---
print("Initializing models...")
# Initialize YOLO segmentation model
model = YOLO("yoloe-11m-seg.pt")
names = ["Access Card", "Screen"]
model.set_classes(names, model.get_text_pe(names))

# Initialize EasyOCR (English only)
reader = easyocr.Reader(['en'], gpu=True)

# Load spaCy model for NER
nlp = spacy.load("en_core_web_trf")

# Initialize Moondream2 for summarization
print("Initializing Moondream2...")
MOONDREAM_LOADED = False
moondream_model = md.vl(model="moondream-2b-int4.mf.gz")

MOONDREAM_LOADED = True
print("All models initialized successfully.")

# Regex patterns for common PIIs + AWS specific
PII_PATTERNS = {
    "EMAIL": r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+",
    "PHONE": r"\+?\d[\d\s-]{11}\d",
    "CREDIT_CARD": r"\b(?:\d[ -]*?){13,16}\b",
    "AWS_ARN": r"arn:aws:[a-z0-9:-]*:[0-9]{12}:[a-zA-Z0-9-_\/]+", # More specific ARN
    "IPv4": r"\b(?:\d{1,3}\s*\.\s*){3}\d{1,3}\b",
    "AWS_SECRET_KEY": r"(?<![A-Z0-9])[A-Za-z0-9/+=]{40}(?![A-Z0-9])" # Common length for AWS secret keys
}

def get_moondream_summary(image_bytes):
    """Generates a summary for an image using Moondream2."""
    if not MOONDREAM_LOADED:
        return "Moondream2 model not loaded, skipping summarization."
    
    try:
        pil_image = Image.open(io.BytesIO(image_bytes)).convert('RGB')
        enc_image = moondream_model.encode_image(pil_image)
        summary = moondream_model.query(
            enc_image, 
            "Describe content of image"
        )
        return summary["answer"]
    except Exception as e:
        return f"Error during Moondream summarization: {e}"

def redact_pii_from_text(text):
    """Redacts PII from a given string of text."""
    redacted_text = text
    redacted_entities = []

    # Regex-based redaction
    for label, pattern in PII_PATTERNS.items():
        # Use re.sub for safer replacement to avoid replacing substrings of already found PII
        redacted_text, count = re.subn(pattern, f"[{label}]", redacted_text)
        if count > 0:
            # Note: This simplification doesn't capture the exact matched text, but logs the action.
             redacted_entities.extend([("redacted", label)] * count)

    # spaCy NER-based redaction
    doc = nlp(redacted_text) 
    temp_text = list(redacted_text)
    for ent in reversed(doc.ents): # Iterate backwards to avoid index shifting
        if ent.label_ in ["PERSON", "GPE", "ORG"]:
            redacted_entities.append((ent.text, ent.label_))
            temp_text[ent.start_char:ent.end_char] = f"[{ent.label_}]"
    redacted_text = "".join(temp_text)
            
    return redacted_text, redacted_entities

def process_image(img_path_or_np, output_path=None):
    """
    Processes a single image file or a numpy array for object and PII redaction.
    Returns the processed image as a numpy array and a summary string.
    """
    if isinstance(img_path_or_np, str):
        img = cv2.imread(img_path_or_np)
        if img is None:
            print(f"Error: Could not read image at {img_path_or_np}")
            return None, ""
    else:
        img = img_path_or_np

    processed_img = img.copy()
    detected_labels = []
    redacted_texts_summary = []

    results = model.predict(processed_img)
    if results[0].masks is not None and results[0].boxes is not None:
        orig_h, orig_w = processed_img.shape[:2]
        for mask_data, box_data in zip(results[0].masks.data.cpu().numpy(), results[0].boxes):
            class_id = int(box_data.cls[0].cpu().numpy())
            class_name = model.names[class_id]

            if class_name == "Access Card":
                detected_labels.append(class_name)
                resized_mask = cv2.resize(mask_data, (orig_w, orig_h), interpolation=cv2.INTER_NEAREST)
                mask_3d_bool = np.stack([resized_mask] * 3, axis=-1).astype(bool)
                processed_img = np.where(mask_3d_bool, 0, processed_img) # Blackout
                
                x1, y1, x2, y2 = map(int, box_data.xyxy[0].cpu().numpy())
                center_x, center_y = int((x1 + x2) / 2), int((y1 + y2) / 2)
                cv2.putText(processed_img, class_name, (center_x - 50, center_y), cv2.FONT_HERSHEY_DUPLEX, 0.9, (255, 255, 255), 2)

    ocr_results = reader.readtext(processed_img)
    skip_next = False
    for (bbox, text, prob) in ocr_results:
        if skip_next or "arn" in text.lower() or "roleid" in text.lower():
            pts = np.array(bbox, dtype=np.int32)
            cv2.fillPoly(processed_img, [pts], (0, 0, 0))
            redacted_texts_summary.append((text, "AWS_IDENTIFIER"))
            skip_next = "arn" in text.lower() or "roleid" in text.lower()
            continue
        
        _, redacted_items = redact_pii_from_text(text)
        if redacted_items:
            pts = np.array(bbox, dtype=np.int32)
            cv2.fillPoly(processed_img, [pts], (0, 0, 0))
            redacted_texts_summary.extend(redacted_items)

    if output_path:
        cv2.imwrite(output_path, processed_img)

    final_text = ""
    if detected_labels:
        final_text += "Detections: " + ", ".join(list(set(detected_labels)))
    if redacted_texts_summary:
        final_text += "\nRedacted Entities: " + str(redacted_texts_summary)
    
    return processed_img, final_text

def process_pdf(input_path, output_path):
    """Redacts PII from text content in a PDF file."""
    doc = fitz.open(input_path)
    redacted_count = 0
    for page in doc:
        text = page.get_text("text")
        _, redacted_entities = redact_pii_from_text(text)
        
        pii_texts_to_redact = [item[0] for item in redacted_entities if item[0] != 'redacted']

        for text_to_redact in pii_texts_to_redact:
            areas = page.search_for(text_to_redact)
            for inst in areas:
                page.add_redact_annot(inst, fill=(0, 0, 0))
                redacted_count += 1
        
        page.apply_redactions()
        
    doc.save(output_path)
    doc.close()
    return f"Processed PDF. Redacted {redacted_count} text instances."

def process_tabular(input_path, output_path):
    """Redacts PII from CSV or XLSX files."""
    if input_path.endswith('.csv'):
        df = pd.read_csv(input_path, dtype=str)
    else:
        df = pd.read_excel(input_path, dtype=str)
    
    df.fillna('', inplace=True)
    
    redacted_df = df.applymap(lambda x: redact_pii_from_text(str(x))[0] if isinstance(x, str) else x)
    
    if output_path.endswith('.csv'):
        redacted_df.to_csv(output_path, index=False)
    else:
        redacted_df.to_excel(output_path, index=False)
    return f"Processed tabular file. Saved to {output_path}"
    
def process_pptx(input_path, output_path):
    """Redacts PII from text and images in a PowerPoint file."""
    prs = Presentation(input_path)
    total_redactions = 0
    
    for i, slide in enumerate(prs.slides):
        print(f"  Processing slide {i+1}/{len(prs.slides)}")
        shapes_to_remove = []
        new_shapes_to_add = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        redacted_text, entities = redact_pii_from_text(original_text)
                        if entities:
                            run.text = redacted_text
                            total_redactions += len(entities)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        original_text = cell.text
                        redacted_text, entities = redact_pii_from_text(original_text)
                        if entities:
                            cell.text = redacted_text
                            total_redactions += len(entities)
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                pil_img = Image.open(io.BytesIO(image_bytes))
                np_img = np.array(pil_img.convert('RGB')) 
                np_img = cv2.cvtColor(np_img, cv2.COLOR_RGB2BGR)
                redacted_img_np, _ = process_image(np_img)
                _, buffer = cv2.imencode('.png', redacted_img_np)
                io_buf = io.BytesIO(buffer)
                shapes_to_remove.append(shape)
                new_shapes_to_add.append({
                    "left": shape.left, "top": shape.top, 
                    "width": shape.width, "height": shape.height,
                    "image": io_buf
                })

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        for shape_info in new_shapes_to_add:
            slide.shapes.add_picture(
                shape_info["image"], shape_info["left"], shape_info["top"],
                shape_info["width"], shape_info["height"]
            )

    prs.save(output_path)
    return f"Processed PPTX. Performed {total_redactions} text redactions and processed images."

def process_file(input_path, output_dir="output"):
    """
    Identifies file type, calls the appropriate processing function, and summarizes the output.
    """
    if not os.path.exists(input_path):
        print(f"Error: File not found at {input_path}")
        return

    os.makedirs(output_dir, exist_ok=True)
    
    filename = os.path.basename(input_path)
    output_path = os.path.join(output_dir, filename)
    
    print(f"\nProcessing {filename}...")
    
    ext = os.path.splitext(filename)[1].lower()
    summary = ""
    moondream_summary = ""
    
    if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
        _, summary = process_image(input_path, output_path)
        if MOONDREAM_LOADED and os.path.exists(output_path):
            with open(output_path, "rb") as f:
                image_bytes = f.read()
            print("  Generating Moondream2 summary for the image...")
            moondream_summary = get_moondream_summary(image_bytes)

    elif ext == '.pdf':
        summary = process_pdf(input_path, output_path)
        if MOONDREAM_LOADED and os.path.exists(output_path):
            try:
                doc = fitz.open(output_path)
                if len(doc) > 0:
                    page = doc.load_page(0)  # First page
                    pix = page.get_pixmap(dpi=150)
                    image_bytes = pix.tobytes("png")
                    print("  Generating Moondream2 summary for the first page of the PDF...")
                    moondream_summary = get_moondream_summary(image_bytes)
                doc.close()
            except Exception as e:
                moondream_summary = f"Could not extract image from PDF for summarization: {e}"
                
    elif ext in ['.csv', '.xlsx']:
        summary = process_tabular(input_path, output_path)
        # No visual summary for tabular data

    elif ext == '.pptx':
        summary = process_pptx(input_path, output_path)
        if MOONDREAM_LOADED and os.path.exists(output_path):
            try:
                prs = Presentation(output_path)
                image_bytes = None
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_bytes = shape.image.blob
                            break
                    if image_bytes:
                        break
                
                if image_bytes:
                    print("  Generating Moondream2 summary for the first image in the PPTX...")
                    moondream_summary = get_moondream_summary(image_bytes)
                else:
                    moondream_summary = "No images found in the PPTX to summarize."
            except Exception as e:
                moondream_summary = f"Could not extract image from PPTX for summarization: {e}"

    else:
        summary = f"File type '{ext}' is not supported."
        
    print(f"Finished processing {filename}.")
    print(f"Redaction Summary: {summary if summary else 'No redactions performed.'}")
    if moondream_summary:
        print(f"Moondream2 Content Summary: {moondream_summary}")
        
if __name__ == "__main__":
    
    # Process a single file example
    process_file("Files/File_001.png")
    
    # Process all files in a directory
    #file_directory = "Files"
    #for filename in os.listdir(file_directory):
    #    file_path = os.path.join(file_directory, filename)
    #    if os.path.isfile(file_path):
    #        process_file(file_path)